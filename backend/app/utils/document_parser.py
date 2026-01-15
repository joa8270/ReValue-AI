"""
文件解析模組
支援 Word (.docx), PowerPoint (.pptx), TXT 格式解析
"""
import io
from typing import Optional

def parse_docx(file_bytes: bytes) -> str:
    """
    解析 Word (.docx) 文件，提取純文字內容
    
    Args:
        file_bytes: Word 文件的二進位內容
        
    Returns:
        提取的純文字內容
    """
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        
        # 提取所有段落文字
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        # 同時提取表格中的文字
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    table_texts.append(row_text)
        
        full_text = "\n".join(paragraphs)
        if table_texts:
            full_text += "\n\n[表格內容]\n" + "\n".join(table_texts)
            
        return full_text
    except Exception as e:
        return f"[Word 解析失敗: {str(e)}]"


def parse_pptx(file_bytes: bytes) -> str:
    """
    解析 PowerPoint (.pptx) 文件，提取所有投影片的文字內容
    
    Args:
        file_bytes: PPT 文件的二進位內容
        
    Returns:
        提取的純文字內容（按投影片編號排列）
    """
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        
        slides_content = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                slides_content.append(f"[投影片 {i}]\n" + "\n".join(slide_texts))
        
        return "\n\n".join(slides_content)
    except Exception as e:
        return f"[PPT 解析失敗: {str(e)}]"


def parse_txt(file_bytes: bytes) -> str:
    """
    解析純文字 (.txt) 文件
    
    Args:
        file_bytes: 文字文件的二進位內容
        
    Returns:
        文字內容
    """
    try:
        # 嘗試 UTF-8 編碼
        try:
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            pass
        
        # 嘗試 Big5 編碼 (繁體中文常用)
        try:
            return file_bytes.decode('big5')
        except UnicodeDecodeError:
            pass
        
        # 嘗試 GBK 編碼 (簡體中文常用)
        try:
            return file_bytes.decode('gbk')
        except UnicodeDecodeError:
            pass
        
        # 最後使用 latin-1 (不會失敗)
        return file_bytes.decode('latin-1')
    except Exception as e:
        return f"[TXT 解析失敗: {str(e)}]"


def get_file_extension(filename: str) -> str:
    """取得檔案副檔名（小寫）"""
    if not filename:
        return ""
    parts = filename.rsplit(".", 1)
    return parts[-1].lower() if len(parts) > 1 else ""


def parse_document(file_bytes: bytes, filename: str) -> Optional[str]:
    """
    根據檔案類型自動選擇解析器
    
    Args:
        file_bytes: 文件的二進位內容
        filename: 原始檔案名稱
        
    Returns:
        解析後的文字內容，若格式不支援則返回 None
    """
    ext = get_file_extension(filename)
    
    if ext == "docx":
        return parse_docx(file_bytes)
    elif ext == "pptx":
        return parse_pptx(file_bytes)
    elif ext == "txt":
        return parse_txt(file_bytes)
    elif ext == "pdf":
        # PDF 保持原有流程（直接傳給 Gemini 處理）
        return None
    else:
        return None
